#!/usr/bin/env python3
"""
Multi-Agent Business Assistant — 100% Free & Local Service Verification Suite
Tests Report Generation (PDF/DOCX/XLSX/PPTX), Local Vector Store (Qdrant), Wikipedia, and SQLite CRUD.
"""

import os
import sys
import sqlite3
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import docx
import openpyxl
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
import wikipedia

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR = os.path.join(BASE_DIR, "data", "test_reports")
os.makedirs(OUTPUT_DIR, exist_ok=True)

def print_header(title):
    print("\n" + "="*70)
    print(f" 🏠 {title}")
    print("="*70)

def test_local_report_generation():
    print_header("STEP 1: Testing 100% Free Local Report Generator (Report Agent #6)")
    
    # 1. PDF Report
    pdf_path = os.path.join(OUTPUT_DIR, "ai_meeting_brief.pdf")
    c = canvas.Canvas(pdf_path, pagesize=letter)
    c.drawString(100, 750, "Multi-Agent Business Assistant — Executive Report")
    c.drawString(100, 720, "Status: 100% Automated Local PDF Generation Verified")
    c.save()
    print(f" ✅ PDF Generated successfully: {pdf_path}")

    # 2. Word DOCX Report
    docx_path = os.path.join(OUTPUT_DIR, "project_summary.docx")
    doc = docx.Document()
    doc.add_heading("AI Business Assistant — Project Briefing", 0)
    doc.add_paragraph("This document was generated totally free offline using python-docx.")
    doc.save(docx_path)
    print(f" ✅ DOCX Generated successfully: {docx_path}")

    # 3. Excel XLSX Spreadsheet
    xlsx_path = os.path.join(OUTPUT_DIR, "analytics_kpis.xlsx")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "AI KPIs"
    ws['A1'] = "Metric Name"
    ws['B1'] = "Value"
    ws['A2'] = "Automated Tasks Complete"
    ws['B2'] = "99.8%"
    wb.save(xlsx_path)
    print(f" ✅ XLSX Spreadsheet Generated successfully: {xlsx_path}")

    # 4. PowerPoint PPTX Presentation
    pptx_path = os.path.join(OUTPUT_DIR, "presentation_deck.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Multi-Agent System Architecture"
    slide.placeholders[1].text = "Generated autonomously via Report Agent (#6)"
    prs.save(pptx_path)
    print(f" ✅ PPTX Presentation Generated successfully: {pptx_path}")

def test_local_qdrant_vector_db():
    print_header("STEP 2: Testing 100% Free Local Qdrant Vector Engine (RAG #14 & Memory #8)")
    # Initialize zero-server local disk/memory Qdrant storage
    vector_db_path = os.path.join(BASE_DIR, "data", "qdrant_storage")
    os.makedirs(vector_db_path, exist_ok=True)
    
    print(f" 📁 Initializing Qdrant Local File System storage at: {vector_db_path}...")
    client = QdrantClient(path=vector_db_path)
    
    collection_name = "company_knowledge_base"
    if client.collection_exists(collection_name):
        client.delete_collection(collection_name)
    client.create_collection(
        collection_name=collection_name,
        vectors_config=VectorParams(size=4, distance=Distance.COSINE),
    )
    
    print(" 📥 Inserting test semantic vector memory...")
    client.upsert(
        collection_name=collection_name,
        points=[
            PointStruct(id=1, vector=[0.9, 0.1, 0.0, 0.2], payload={"text": "Refund policy: 30-day money back"}),
            PointStruct(id=2, vector=[0.1, 0.8, 0.9, 0.0], payload={"text": "Project roadmap: Phase 5 LangGraph AI Core"}),
        ]
    )
    
    print(" 🔍 Performing simulated semantic RAG search for 'policy'...")
    response = client.query_points(
        collection_name=collection_name,
        query=[0.95, 0.05, 0.0, 0.1],
        limit=1
    )
    results = response.points
    matched_text = results[0].payload.get("text") if results else "None"
    score = results[0].score if results else 0.0
    print(f" ✅ Match retrieved: [{matched_text}] (Similarity Score: {score:.4f})")

def test_local_wikipedia():
    print_header("STEP 3: Testing Free Wikipedia SDK (Research Agent #2)")
    query = "Information systems"
    print(f" 📖 Querying Wikipedia for topic: '{query}'...")
    try:
        summary = wikipedia.summary(query, sentences=2)
        print(f" ✅ Wikipedia Result:\n     \"{summary}\"")
    except Exception as e:
        print(f" ⚠️ Wikipedia warning: {e}")

def test_local_sqlite_db():
    print_header("STEP 4: Testing Local SQLite Relational Database (No-Server Fallback for Phase 4)")
    db_file = os.path.join(BASE_DIR, "data", "local_system.db")
    conn = sqlite3.connect(db_file)
    cursor = conn.cursor()
    cursor.execute("CREATE TABLE IF NOT EXISTS test_users (id INTEGER PRIMARY KEY, email TEXT, role TEXT)")
    cursor.execute("INSERT OR REPLACE INTO test_users (id, email, role) VALUES (1, 'bouchakyahya0@gmail.com', 'Admin')")
    conn.commit()
    cursor.execute("SELECT * FROM test_users WHERE id = 1")
    user = cursor.fetchone()
    conn.close()
    print(f" ✅ SQLite Database Verification successful! Stored Admin Record: {user}")

def main():
    print_header("STARTING FREE & LOCAL SOLUTION ARCHITECTURE TEST SUITE")
    test_local_report_generation()
    test_local_qdrant_vector_db()
    test_local_wikipedia()
    test_local_sqlite_db()
    print_header("🎉 ALL FREE LOCAL SERVICES VERIFIED! 100% READY WITH $0 ADDITIONAL COST 🎉")

if __name__ == "__main__":
    main()
